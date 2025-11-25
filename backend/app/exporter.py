from docx import Document
from pptx import Presentation
from io import BytesIO

def export_docx(project, sections):
    doc = Document()
    doc.add_heading(project.title, level=1)
    if project.prompt:
        doc.add_paragraph(f"Prompt: {project.prompt}")
    for sec in sections:
        doc.add_heading(sec.title, level=2)
        doc.add_paragraph(sec.content or '')
    bio = BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio


def export_pptx(project, sections):
    prs = Presentation()
    # Title slide
    slide_layout = prs.slide_layouts[5]
    for sec in sections:
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = sec.title
        # add content as textbox
        txBox = slide.shapes.add_textbox(left=1000000, top=1500000, width=8000000, height=3000000)
        tf = txBox.text_frame
        tf.text = sec.content or ''

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio
